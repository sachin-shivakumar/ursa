import json
import os
import re
import shutil
import subprocess
import unicodedata
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path
from typing import Any, Optional
from urllib.parse import urljoin, urlparse

import justext
import requests
import trafilatura
from bs4 import BeautifulSoup
from langchain_community.document_loaders import PyPDFLoader

# Check for optional dependencies
docx_installed = False
pptx_installed = False
try:
    from docx import Document

    docx_installed = True
except Exception:
    pass

try:
    from pptx import Presentation

    pptx_installed = True
except Exception:
    pass


# Curate this for your environment. Start broad, tighten later.
TEXT_EXTENSIONS = {
    # plain text & docs
    ".txt",
    ".md",
    ".rst",
    ".rtf",
    ".tex",
    ".csv",
    ".tsv",
    ".json",
    ".jsonl",
    ".log",
    ".xml",
    ".html",
    ".htm",
    # source code (common)
    ".py",
    ".pyi",
    ".ipynb",
    ".c",
    ".h",
    ".cpp",
    ".hpp",
    ".cc",
    ".java",
    ".kt",
    ".scala",
    ".js",
    ".ts",
    ".jsx",
    ".tsx",
    ".go",
    ".rs",
    ".rb",
    ".php",
    ".sh",
    ".bash",
    ".zsh",
    ".ps1",
}

SPECIAL_TEXT_FILENAMES = {
    "makefile",
    "readme",
    "license",
}

OFFICE_EXTENSIONS = {".docx", ".pptx", ".odt", ".odp"}


def extract_json(text: str) -> list[dict]:
    """
    Extract a JSON object or array from text that might contain markdown or other content.

    The function attempts three strategies:
        1. Extract JSON from a markdown code block labeled as JSON.
        2. Extract JSON from any markdown code block.
        3. Use bracket matching to extract a JSON substring starting with '{' or '['.

    Returns:
        A Python object parsed from the JSON string (dict or list).

    Raises:
        ValueError: If no valid JSON is found.
    """
    # Approach 1: Look for a markdown code block specifically labeled as JSON.
    labeled_block = re.search(
        r"```json\s*([\[{].*?[\]}])\s*```", text, re.DOTALL
    )
    if labeled_block:
        json_str = labeled_block.group(1).strip()
        try:
            return json.loads(json_str)
        except json.JSONDecodeError:
            # Fall back to the next approach if parsing fails.
            pass

    # Approach 2: Look for any code block delimited by triple backticks.
    generic_block = re.search(r"```(.*?)```", text, re.DOTALL)
    if generic_block:
        json_str = generic_block.group(1).strip()
        if json_str.startswith("{") or json_str.startswith("["):
            try:
                return json.loads(json_str)
            except json.JSONDecodeError:
                pass

    # Approach 3: Attempt to extract JSON using bracket matching.
    # Find the first occurrence of either '{' or '['.
    first_obj = text.find("{")
    first_arr = text.find("[")
    if first_obj == -1 and first_arr == -1:
        raise ValueError("No JSON object or array found in the text.")

    # Determine which bracket comes first.
    if first_obj == -1:
        start = first_arr
        open_bracket = "["
        close_bracket = "]"
    elif first_arr == -1:
        start = first_obj
        open_bracket = "{"
        close_bracket = "}"
    else:
        if first_obj < first_arr:
            start = first_obj
            open_bracket = "{"
            close_bracket = "}"
        else:
            start = first_arr
            open_bracket = "["
            close_bracket = "]"

    # Bracket matching: find the matching closing bracket.
    depth = 0
    end = None
    for i in range(start, len(text)):
        if text[i] == open_bracket:
            depth += 1
        elif text[i] == close_bracket:
            depth -= 1
            if depth == 0:
                end = i
                break

    if end is None:
        raise ValueError(
            "Could not find matching closing bracket for JSON content."
        )

    json_str = text[start : end + 1]
    try:
        return json.loads(json_str)
    except json.JSONDecodeError as e:
        raise ValueError("Extracted content is not valid JSON.") from e


PDF_CT_HINTS = (
    "application/pdf",
    "binary/octet-stream",
)  # some servers mislabel
PDF_EXT_RE = re.compile(r"\.pdf($|\?)", re.IGNORECASE)


def _is_pdf_response(resp: requests.Response) -> bool:
    ct = resp.headers.get("Content-Type", "").lower()
    if any(hint in ct for hint in PDF_CT_HINTS):
        return True
    # Sometimes servers omit CT but set filename
    cd = resp.headers.get("Content-Disposition", "")
    if "filename" in cd and ".pdf" in cd.lower():
        return True
    # Last resort: URL extension
    return bool(PDF_EXT_RE.search(resp.url))


def _derive_filename_from_cd_or_url(
    resp: requests.Response, fallback: str
) -> str:
    cd = resp.headers.get("Content-Disposition", "")
    m = re.search(r'filename\*?=(?:UTF-8\'\')?"?([^\";]+)"?', cd, re.IGNORECASE)
    if m:
        name = m.group(1)
        # Some headers include quotes
        name = name.strip("\"'")

        # RFC 5987 may encode UTF-8 in filename*; we’re treating as plain here.
        if not name.lower().endswith(".pdf"):
            name += ".pdf"
        return name

    # use URL last path segment if looks like PDF
    parsed = urlparse(resp.url)
    base = os.path.basename(parsed.path) or fallback
    if not base.lower().endswith(".pdf"):
        if PDF_EXT_RE.search(resp.url):
            base = re.sub(
                r"(\.pdf)(?:$|\?).*", r"\1", base, flags=re.IGNORECASE
            )
            if not base.lower().endswith(".pdf"):
                base += ".pdf"
        else:
            base = (
                fallback
                if fallback.lower().endswith(".pdf")
                else fallback + ".pdf"
            )
    return base


def _download_stream_to(path: str, resp: requests.Response) -> str:
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "wb") as f:
        shutil.copyfileobj(resp.raw, f)
    return path


def _get_soup(
    url: str, timeout: int = 20, headers: Optional[dict[str, str]] = None
) -> BeautifulSoup:
    r = requests.get(url, timeout=timeout, headers=headers or {})
    r.raise_for_status()
    return BeautifulSoup(r.text, "html.parser")


def _find_pdf_on_landing(soup: BeautifulSoup, base_url: str) -> Optional[str]:
    # 1) meta citation_pdf_url
    meta = soup.find("meta", attrs={"name": "citation_pdf_url"})
    if meta and meta.get("content"):
        return urljoin(base_url, meta["content"])

    # 2) obvious anchors: text contains 'PDF' or 'Download'
    for a in soup.find_all("a", href=True):
        label = (a.get_text(" ", strip=True) or "").lower()
        href = a["href"]
        if "pdf" in label or "download" in label or PDF_EXT_RE.search(href):
            return urljoin(base_url, href)

    # 3) buttons that wrap an anchor
    for btn in soup.find_all(["button", "a"], href=True):
        label = (btn.get_text(" ", strip=True) or "").lower()
        href = btn.get("href")
        if href and (
            "pdf" in label or "download" in label or PDF_EXT_RE.search(href)
        ):
            return urljoin(base_url, href)

    return None


def _pdf_page_count(path: Path) -> int:
    try:
        loader = PyPDFLoader(path)
        pages = loader.load()
        return len(pages)
    except Exception as e:
        print("[Error]: ", e)
        return 0


def ocrmypdf_is_installed() -> bool:
    return shutil.which("ocrmypdf") is not None


def _ocr_to_searchable_pdf(
    src_pdf: str, out_pdf: str, *, mode: str = "skip"
) -> None:
    # mode:
    #  - "skip":  only OCR pages that look like they need it (your current behavior)
    #  - "force": rasterize + OCR everything (fixes vector/outlined “no images” PDFs)
    if not ocrmypdf_is_installed():
        raise ImportError(
            "ocrmypdf was not found in your path. "
            "See installation instructions:"
            "https://github.com/ocrmypdf/OCRmyPDF?tab=readme-ov-file#installation"
        )

    cmd = ["ocrmypdf", "--rotate-pages", "--deskew", "--clean"]

    if mode == "force":
        cmd += ["--force-ocr"]
    else:
        cmd += ["--skip-text"]

    # Optional: dump a sidecar text file for debugging confidence
    if os.getenv("READ_FILE_OCR_SIDECAR", "0").lower() in ("1", "true", "yes"):
        cmd += ["--sidecar", out_pdf + ".txt"]

    cmd += [src_pdf, out_pdf]

    # Don’t swallow stderr/stdout when debugging
    debug = os.getenv("READ_FILE_OCR_DEBUG", "0").lower() in (
        "1",
        "true",
        "yes",
    )
    subprocess.run(
        cmd,
        check=True,
        stdout=None if debug else subprocess.PIPE,
        stderr=None if debug else subprocess.PIPE,
        text=True,
    )


def resolve_pdf_from_osti_record(
    rec: dict[str, Any],
    *,
    headers: Optional[dict[str, str]] = None,
    unpaywall_email: Optional[str] = None,
    timeout: int = 25,
) -> tuple[Optional[str], Optional[str], str]:
    """
    Returns (pdf_url, landing_used, note)
      - pdf_url: direct downloadable PDF URL if found (or a strong candidate)
      - landing_used: landing page URL we parsed (if any)
      - note: brief trace of how we found it
    """
    headers = headers or {"User-Agent": "Mozilla/5.0"}
    note_parts: list[str] = []

    links = rec.get("links", []) or []
    # doi = rec.get("doi")

    # 1) Try 'fulltext' first (OSTI purl)
    fulltext = None
    for link in links:
        if link.get("rel") == "fulltext":
            fulltext = link.get("href")
            break

    if fulltext:
        note_parts.append("Tried links[fulltext] purl")
        try:
            # Follow redirects; stream to peek headers without loading whole body
            r = requests.get(
                fulltext,
                headers=headers,
                timeout=timeout,
                allow_redirects=True,
                stream=True,
            )
            r.raise_for_status()

            if _is_pdf_response(r):
                note_parts.append("fulltext resolved directly to PDF")
                return (r.url, None, " | ".join(note_parts))

            # Not a PDF: parse page HTML for meta or obvious PDF anchors
            # (If server sent binary but CT lied, _is_pdf_response would have caught via CD or ext)
            r.close()
            soup = _get_soup(fulltext, timeout=timeout, headers=headers)
            candidate = _find_pdf_on_landing(soup, fulltext)
            if candidate:
                note_parts.append(
                    "found PDF via meta/anchor on fulltext landing"
                )
                return (candidate, fulltext, " | ".join(note_parts))
        except Exception as e:
            note_parts.append(f"fulltext failed: {e}")

    # 2) Try DOE PAGES landing (citation_doe_pages)
    doe_pages = None
    for link in links:
        if link.get("rel") == "citation_doe_pages":
            doe_pages = link.get("href")
            break

    if doe_pages:
        note_parts.append("Tried links[citation_doe_pages] landing")
        try:
            soup = _get_soup(doe_pages, timeout=timeout, headers=headers)
            candidate = _find_pdf_on_landing(soup, doe_pages)
            if candidate:
                # Candidate may itself be a landing—check if it serves PDF
                try:
                    r2 = requests.get(
                        candidate,
                        headers=headers,
                        timeout=timeout,
                        allow_redirects=True,
                        stream=True,
                    )
                    r2.raise_for_status()
                    if _is_pdf_response(r2):
                        note_parts.append("citation_doe_pages → direct PDF")
                        return (r2.url, doe_pages, " | ".join(note_parts))
                    r2.close()
                except Exception:
                    pass
                # If not clearly PDF, still return as a candidate (agent will fetch & parse)
                note_parts.append(
                    "citation_doe_pages → PDF-like candidate (not confirmed by headers)"
                )
                return (candidate, doe_pages, " | ".join(note_parts))
        except Exception as e:
            note_parts.append(f"citation_doe_pages failed: {e}")

    # # 3) Optional: DOI → Unpaywall OA
    # if doi and unpaywall_email:
    #     note_parts.append("Tried Unpaywall via DOI")
    #     pdf_from_ua = _resolve_pdf_via_unpaywall(doi, unpaywall_email)
    #     if pdf_from_ua:
    #         # May be direct PDF or landing; the caller will validate headers during download
    #         note_parts.append("Unpaywall returned candidate")
    #         return (pdf_from_ua, None, " | ".join(note_parts))

    # 4) Give up
    note_parts.append("No PDF found")
    return (None, None, " | ".join(note_parts))


def _normalize_ws(text: str) -> str:
    # Normalize unicode, collapse whitespace, and strip control chars
    text = unicodedata.normalize("NFKC", text)
    text = re.sub(r"[ \t\r\f\v]+", " ", text)
    text = re.sub(r"\s*\n\s*", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = text.strip()
    return text


def _dedupe_lines(text: str, min_len: int = 40) -> str:
    seen = set()
    out = []
    for line in text.splitlines():
        stripped = line.strip()
        # Ignore very short or repeated lines (menus, cookie banners, etc.)
        if len(stripped) < min_len:
            continue
        key = stripped.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(stripped)
    return "\n\n".join(out)


def extract_main_text_only(html: str, *, max_chars: int = 250_000) -> str:
    """
    Returns plain text with navigation/ads/scripts removed.
    Prefers trafilatura -> jusText -> BS4 paragraphs.
    """
    # 1) Trafilatura
    # You can tune config: with_metadata, include_comments, include_images, favor_recall, etc.
    cfg = trafilatura.settings.use_config()
    cfg.set("DEFAULT", "include_comments", "false")
    cfg.set("DEFAULT", "include_tables", "false")
    cfg.set("DEFAULT", "favor_recall", "false")  # be stricter; less noise
    try:
        # If you fetched HTML already, use extract() on string; otherwise, fetch_url(url)
        txt = trafilatura.extract(
            html,
            config=cfg,
            include_comments=False,
            include_tables=False,
            favor_recall=False,
        )
        if txt and txt.strip():
            txt = _normalize_ws(txt)
            txt = _dedupe_lines(txt)
            return txt[:max_chars]
    except Exception:
        pass

    # 2) jusText
    try:
        paragraphs = justext.justext(html, justext.get_stoplist("English"))
        body_paras = [p.text for p in paragraphs if not p.is_boilerplate]
        if body_paras:
            txt = _normalize_ws("\n\n".join(body_paras))
            txt = _dedupe_lines(txt)
            return txt[:max_chars]
    except Exception:
        pass

    # 4) last-resort: BS4 paragraphs/headings only
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup([
        "script",
        "style",
        "noscript",
        "header",
        "footer",
        "nav",
        "form",
        "aside",
    ]):
        tag.decompose()
    chunks = []
    for el in soup.find_all(["h1", "h2", "h3", "p", "li", "figcaption"]):
        t = el.get_text(" ", strip=True)
        if t:
            chunks.append(t)
    txt = _normalize_ws("\n\n".join(chunks))
    txt = _dedupe_lines(txt)
    return txt[:max_chars]


def read_text_pdf(path: str | Path) -> str:
    loader = PyPDFLoader(path)
    pages = loader.load()
    return "\n".join(p.page_content for p in pages)


def read_pdf(path: str | Path) -> str:
    full_filename = Path(path)

    try:
        # 1) normal extraction
        text = read_text_pdf(full_filename) or ""

        # 2) decide if OCR fallback is needed
        pages = _pdf_page_count(full_filename)
        ocr_enabled = os.getenv("READ_FILE_OCR", "1").lower() in (
            "1",
            "true",
            "yes",
        )
        min_pages = int(os.getenv("READ_FILE_OCR_MIN_PAGES", "3"))
        min_chars = int(os.getenv("READ_FILE_OCR_MIN_CHARS", "3000"))

        if ocr_enabled and pages >= min_pages and len(text) < min_chars:
            src = Path(full_filename)

            mode_env = os.getenv("READ_FILE_OCR_MODE", "auto").lower()
            force_if_still_low = os.getenv(
                "READ_FILE_OCR_FORCE_IF_STILL_LOW", "1"
            ).lower() in ("1", "true", "yes")

            try:
                # First pass (skip-text) unless user forces always-force
                first_mode = "force" if mode_env == "force" else "skip"
                ocr_pdf = str(
                    src.with_suffix(src.suffix + f".ocr.{first_mode}.pdf")
                )

                if not os.path.exists(ocr_pdf) or os.path.getmtime(
                    ocr_pdf
                ) < os.path.getmtime(full_filename):
                    print(
                        f"[OCR]: mode={first_mode} ({len(text)} chars, {pages} pages) -> {ocr_pdf}"
                    )
                    _ocr_to_searchable_pdf(
                        full_filename, ocr_pdf, mode=first_mode
                    )
                else:
                    print(f"[OCR]: using cached OCR PDF -> {ocr_pdf}")

                text2 = read_text_pdf(ocr_pdf) or ""
                if len(text2) > len(text):
                    text = text2

                # Second pass: if still low and we weren’t already forcing, try force-ocr
                if (
                    force_if_still_low
                    and mode_env != "force"
                    and len(text) < min_chars
                ):
                    force_pdf = str(
                        src.with_suffix(src.suffix + ".ocr.force.pdf")
                    )
                    if not os.path.exists(force_pdf) or os.path.getmtime(
                        force_pdf
                    ) < os.path.getmtime(full_filename):
                        print(
                            f"[OCR]: still low after skip-text; retrying with force-ocr -> {force_pdf}"
                        )
                        _ocr_to_searchable_pdf(
                            full_filename, force_pdf, mode="force"
                        )
                    else:
                        print(
                            f"[OCR]: using cached force OCR PDF -> {force_pdf}"
                        )

                    text3 = read_text_pdf(force_pdf) or ""
                    if len(text3) > len(text):
                        text = text3

            except (FileNotFoundError, subprocess.CalledProcessError) as e:
                # Missing ocrmypdf or OCR failed: keep original extraction
                print(f"[OCR Error]: {e}")
            except Exception as e:
                # Any other OCR-related failure: keep original extraction
                print(f"[OCR Error]: {e}")

        return text

    except subprocess.CalledProcessError as e:
        # OCR failed; return whatever we got from normal extraction
        err = (e.stderr or "")[:500]
        print(f"[OCR Error]: {err}")
        return text if text else f"[Error]: OCR failed: {err}"
    except Exception as e:
        print(f"[Error]: {e}")
        return f"[Error]: {e}"


def read_text_file(path: str | Path) -> str:
    """
    Reads in a file at a given path into a string

    Args:
        path: string filename, with path, to read in
    """
    with open(path, "r", encoding="utf-8") as file:
        file_contents = file.read()
    return file_contents


# helper to extract text from OpenDocument formats (.odt/.odp)
def read_odf(path: Path) -> str:
    with zipfile.ZipFile(path, "r") as zf:
        xml_bytes = zf.read("content.xml")

    root = ET.fromstring(xml_bytes)
    # simple, robust extraction: gather all text nodes
    chunks = [t.strip() for t in root.itertext() if t and t.strip()]
    return "\n".join(chunks)


# helper to parse .docx via python-docx
def read_docx(path: Path) -> str:
    if docx_installed:
        doc = Document(str(path))
        parts: list[str] = []

        for para in doc.paragraphs:
            txt = (para.text or "").strip()
            if txt:
                parts.append(txt)

        # also pull table text
        for table in doc.tables:
            for row in table.rows:
                row_txt = "\t".join(
                    (cell.text or "").strip() for cell in row.cells
                ).strip()
                if row_txt:
                    parts.append(row_txt)

        return "\n".join(parts)
    else:
        return (
            f"No DOCX reader so skipping {str(path)}.\n",
            "Consider installing via `pip install 'ursa-ai[office_readers]'`.",
        )


# helper to parse .pptx via python-pptx
def read_pptx(path: Path) -> str:
    if pptx_installed:
        prs = Presentation(str(path))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = (shape.text or "").strip()
                    if txt:
                        parts.append(txt)
        return "\n".join(parts)
    else:
        return (
            f"No PPTX reader so skipping {str(path)}.\n",
            "Consider installing via `pip install 'ursa-ai[office_readers]'`.",
        )


def read_text_from_file(path):
    custom_extensions = [
        item.strip()
        for item in os.environ.get("URSA_TEXT_EXTENSIONS", "").split(",")
    ]
    custom_readable_files = [
        item.strip()
        for item in os.environ.get("URSA_SPECIAL_TEXT_FILENAMES", "").split(",")
    ]
    ext = path.suffix.lower()
    try:
        match ext:
            case ".pdf":
                full_text = read_pdf(path)
            case ".odt" | ".odp":
                full_text = read_odf(path)
            case ".docx":
                full_text = read_docx(path)
            case ".pptx":
                full_text = read_pptx(path)
            case _:
                if (
                    ext in TEXT_EXTENSIONS
                    or path.name.lower() in SPECIAL_TEXT_FILENAMES
                    or ext in custom_extensions
                    or path.name.lower() in custom_readable_files
                ):
                    full_text = read_text_file(path)
                else:
                    full_text = f"Unsupported file type: {path.name}"
    except Exception as e:
        full_text = f"Error loading {path.name}: {e}"
    return full_text
